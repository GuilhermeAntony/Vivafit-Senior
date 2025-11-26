#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para gerar apresentação PowerPoint do TCC VivaFit Seniors
Baseado no modelo de apresentação perante banca
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def criar_apresentacao():
    """Cria apresentação PowerPoint do TCC"""
    
    # Criar apresentação
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Definir cores do tema
    COR_PRIMARIA = RGBColor(14, 165, 163)  # #0ea5a3 (teal)
    COR_TEXTO = RGBColor(51, 51, 51)
    COR_SECUNDARIA = RGBColor(102, 102, 102)
    
    # ==================== SLIDE 1: CAPA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "VIVAFIT SENIORS"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = COR_PRIMARIA
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Aplicativo Mobile de Fitness para Idosos com\nArquitetura Offline-First"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = COR_TEXTO
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Autor e data
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.text = "Guilherme Antony\nTrabalho de Conclusão de Curso\nNovembro de 2025"
    for p in info_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = COR_SECUNDARIA
        p.alignment = PP_ALIGN.CENTER
    
    # ==================== SLIDE 2: INTRODUÇÃO ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "INTRODUÇÃO"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Objetivo Geral
    p = tf.paragraphs[0]
    p.text = "OBJETIVO GERAL"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Desenvolver um aplicativo mobile de exercícios físicos para o público idoso, implementando arquitetura em camadas com padrão offline-first para garantir disponibilidade contínua."
    p.font.size = Pt(18)
    p.font.color.rgb = COR_TEXTO
    p.space_after = Pt(20)
    p.level = 0
    
    # Objetivos Específicos
    p = tf.add_paragraph()
    p.text = "OBJETIVOS ESPECÍFICOS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    objetivos = [
        "Implementar arquitetura em 4 camadas com separação de responsabilidades",
        "Desenvolver sistema de cache offline com taxa de acerto ~85%",
        "Criar catálogo de 10+ exercícios em 4 categorias",
        "Implementar autenticação segura com OAuth 2.0 e JWT"
    ]
    
    for obj in objetivos:
        p = tf.add_paragraph()
        p.text = f"• {obj}"
        p.font.size = Pt(16)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    # ==================== SLIDE 3: PROBLEMA E JUSTIFICATIVA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "PROBLEMA DE PESQUISA E JUSTIFICATIVA"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Problema
    p = tf.paragraphs[0]
    p.text = "PROBLEMA DE PESQUISA"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Como desenvolver um aplicativo mobile de fitness para idosos que mantenha funcionalidade mesmo sem conexão à internet, garantindo disponibilidade contínua dos recursos essenciais?"
    p.font.size = Pt(18)
    p.font.color.rgb = COR_TEXTO
    p.space_after = Pt(20)
    
    # Justificativa
    p = tf.add_paragraph()
    p.text = "JUSTIFICATIVA"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    justificativas = [
        "Envelhecimento populacional crescente (ONU, 2019)",
        "Necessidade de soluções tecnológicas acessíveis para idosos",
        "Falhas de conectividade não devem impedir exercícios físicos",
        "Conformidade com LGPD para dados sensíveis de saúde",
        "Gap no mercado de apps fitness focados no público sênior"
    ]
    
    for just in justificativas:
        p = tf.add_paragraph()
        p.text = f"• {just}"
        p.font.size = Pt(16)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    # ==================== SLIDE 4: FUNDAMENTAÇÃO TEÓRICA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "FUNDAMENTAÇÃO TEÓRICA"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    topicos = [
        ("Arquitetura em Camadas", "Fowler (2002), Bass et al. (2012)", 
         "Separação de responsabilidades, manutenibilidade e escalabilidade"),
        
        ("Padrão Offline-First", "Firtman (2016), Taivalsaari & Mikkonen (2021)",
         "Prioriza funcionamento local, sincronização em background"),
        
        ("React Native e Expo", "Facebook (2023), Expo (2023)",
         "Desenvolvimento multiplataforma, hot reload, APIs nativas"),
        
        ("Autenticação e Segurança", "Stallings & Brown (2018), Hardt (2012)",
         "OAuth 2.0, JWT, Row Level Security (RLS)"),
        
        ("Backend-as-a-Service", "Supabase (2023)",
         "PostgreSQL, autenticação, storage, APIs em tempo real")
    ]
    
    for topico, autores, desc in topicos:
        p = tf.paragraphs[0] if topico == topicos[0][0] else tf.add_paragraph()
        p.text = f"{topico}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COR_PRIMARIA
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = f"Autores: {autores}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = COR_SECUNDARIA
        p.level = 1
        p.space_after = Pt(2)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(12)
    
    # ==================== SLIDE 5: METODOLOGIA - PARTE 1 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "METODOLOGIA"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Tipo de Pesquisa
    p = tf.paragraphs[0]
    p.text = "TIPO DE PESQUISA"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Exploratória e Aplicada - Desenvolvimento de solução tecnológica com análise de requisitos e implementação prática"
    p.font.size = Pt(16)
    p.font.color.rgb = COR_TEXTO
    p.space_after = Pt(16)
    
    # Coleta de Dados
    p = tf.add_paragraph()
    p.text = "COLETA DE DADOS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(8)
    
    coleta = [
        "Pesquisa bibliográfica sobre arquitetura de software mobile",
        "Análise de aplicativos fitness existentes no mercado",
        "Levantamento de requisitos focados no público idoso",
        "Documentação técnica de frameworks e bibliotecas"
    ]
    
    for item in coleta:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(8)
    
    # Unidade de Análise
    p = tf.add_paragraph()
    p.text = "UNIDADE DE ANÁLISE"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Aplicativo VivaFit Seniors desenvolvido com React Native (Expo SDK 54), TypeScript e Supabase"
    p.font.size = Pt(16)
    p.font.color.rgb = COR_TEXTO
    
    # ==================== SLIDE 6: METODOLOGIA - PARTE 2 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "METODOLOGIA (continuação)"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Amostra
    p = tf.paragraphs[0]
    p.text = "AMOSTRA"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Público-alvo: Idosos interessados em manter atividade física regular\nFuncionalidades testadas: 6 telas principais, 10+ exercícios, sistema de cache offline"
    p.font.size = Pt(16)
    p.font.color.rgb = COR_TEXTO
    p.space_after = Pt(20)
    
    # Abordagem
    p = tf.add_paragraph()
    p.text = "ABORDAGEM"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Qualitativa e Quantitativa"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COR_TEXTO
    p.space_after = Pt(8)
    
    abordagens = [
        "Qualitativa: Análise de usabilidade, acessibilidade e experiência do usuário",
        "Quantitativa: Métricas de performance (taxa de cache ~85%, redução de latência 73%)",
        "Metodologia ágil com iterações incrementais",
        "Testes de integração entre camadas da arquitetura"
    ]
    
    for item in abordagens:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    # ==================== SLIDE 7: ARQUITETURA DO SISTEMA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "ARQUITETURA DO SISTEMA"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Arquitetura em 4 Camadas"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)
    
    camadas = [
        ("1. PRESENTATION LAYER (Apresentação)", 
         "Screens (6 telas principais), UI Components reutilizáveis, Verificação de autenticação"),
        
        ("2. BUSINESS LOGIC LAYER (Lógica de Negócio)",
         "Hooks personalizados, Validações, Gerenciamento de estado"),
        
        ("3. DATA ACCESS LAYER (Acesso a Dados)",
         "Cliente Supabase, Sistema de cache offline, AsyncStorage + FileSystem"),
        
        ("4. INFRASTRUCTURE LAYER (Infraestrutura)",
         "React Navigation, Design tokens, Configurações globais")
    ]
    
    for camada, descricao in camadas:
        p = tf.add_paragraph()
        p.text = camada
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COR_PRIMARIA
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = descricao
        p.font.size = Pt(14)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(10)
    
    # ==================== SLIDE 8: SISTEMA DE CACHE OFFLINE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "SISTEMA DE CACHE OFFLINE-FIRST"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ESTRATÉGIA IMPLEMENTADA"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    estrategias = [
        "Política de expiração: 7 dias para dados e imagens",
        "Verificação de validade temporal antes de servir cache",
        "Armazenamento local: AsyncStorage (JSON) + FileSystem (imagens)",
        "Padrão stale-while-revalidate: serve cache e atualiza em background",
        "Limpeza automática de arquivos expirados"
    ]
    
    for est in estrategias:
        p = tf.add_paragraph()
        p.text = f"• {est}"
        p.font.size = Pt(16)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "RESULTADOS DE PERFORMANCE"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    resultados = [
        "Taxa de acerto de cache: ~85%",
        "Redução de tempo de carregamento: 73%",
        "Funcionamento completo offline após primeiro acesso",
        "Experiência de usuário fluida e consistente"
    ]
    
    for res in resultados:
        p = tf.add_paragraph()
        p.text = f"✓ {res}"
        p.font.size = Pt(16)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    # ==================== SLIDE 9: AUTENTICAÇÃO E SEGURANÇA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "AUTENTICAÇÃO E SEGURANÇA"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ARQUITETURA MULTICAMADAS DE SEGURANÇA"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    componentes = [
        ("OAuth 2.0", "Autenticação via Google com fluxo seguro de autorização"),
        ("JWT (JSON Web Tokens)", "Tokens assinados para validação stateless de sessões"),
        ("Row Level Security (RLS)", "Isolamento completo de dados por usuário no PostgreSQL"),
        ("HTTPS/TLS", "Criptografia em trânsito para todas comunicações"),
        ("Conformidade LGPD", "Proteção de dados sensíveis de saúde")
    ]
    
    for comp, desc in componentes:
        p = tf.add_paragraph()
        p.text = comp
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COR_PRIMARIA
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Fluxo: Usuário → App React Native → Supabase Auth → PostgreSQL com RLS → Resposta Protegida"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COR_SECUNDARIA
    p.alignment = PP_ALIGN.CENTER
    
    # ==================== SLIDE 10: FUNCIONALIDADES IMPLEMENTADAS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "FUNCIONALIDADES IMPLEMENTADAS"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    funcionalidades = [
        ("📱 6 Telas Principais", 
         "Dashboard, Perfil, Exercícios, Treino, Progresso, Histórico"),
        
        ("💪 Catálogo de Exercícios",
         "10+ exercícios em 4 categorias: Cardio, Força, Flexibilidade, Equilíbrio"),
        
        ("📊 Acompanhamento de Progresso",
         "Histórico detalhado, estatísticas, gráficos de evolução"),
        
        ("🎯 Planos Personalizados",
         "Treinos adaptados ao nível de condicionamento físico"),
        
        ("👤 Perfil Completo",
         "Informações de saúde, preferências, metas de atividade"),
        
        ("📴 Modo Offline",
         "Acesso completo aos exercícios sem internet")
    ]
    
    for titulo, desc in funcionalidades:
        p = tf.paragraphs[0] if titulo == funcionalidades[0][0] else tf.add_paragraph()
        p.text = titulo
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COR_PRIMARIA
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(10)
    
    # ==================== SLIDE 11: RESULTADOS E DISCUSSÃO ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "APRESENTAÇÃO E DISCUSSÃO DOS RESULTADOS"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MÉTRICAS DE SUCESSO"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    metricas = [
        "✓ Taxa de acerto de cache: ~85% (otimização significativa)",
        "✓ Redução de latência: 73% vs requisições diretas",
        "✓ Funcionamento offline completo após primeira sincronização",
        "✓ Arquitetura escalável e manutenível em 4 camadas",
        "✓ Type safety com TypeScript (zero erros de tipo em produção)",
        "✓ Autenticação segura com OAuth 2.0 + JWT + RLS"
    ]
    
    for metrica in metricas:
        p = tf.add_paragraph()
        p.text = metrica
        p.font.size = Pt(15)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "BENEFÍCIOS ALCANÇADOS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    beneficios = [
        "Experiência de usuário fluida e consistente",
        "Separação clara de responsabilidades facilita manutenção",
        "Base sólida para evolução futura do sistema",
        "Conformidade com LGPD para dados de saúde"
    ]
    
    for ben in beneficios:
        p = tf.add_paragraph()
        p.text = f"• {ben}"
        p.font.size = Pt(15)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    # ==================== SLIDE 12: TECNOLOGIAS UTILIZADAS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "TECNOLOGIAS UTILIZADAS"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    categorias = [
        ("Frontend Mobile", [
            "React Native 0.76.x",
            "Expo SDK 54.0.0",
            "TypeScript 5.3.x",
            "React Navigation 6.x"
        ]),
        ("Backend e Autenticação", [
            "Supabase (PostgreSQL)",
            "OAuth 2.0 / JWT",
            "Row Level Security (RLS)"
        ]),
        ("Storage e Cache", [
            "AsyncStorage 1.23.x",
            "Expo FileSystem 17.x",
            "Cache offline de 7 dias"
        ]),
        ("Build e Deploy", [
            "EAS (Expo Application Services)",
            "APK para Android 5.0+"
        ])
    ]
    
    for cat, techs in categorias:
        p = tf.paragraphs[0] if cat == categorias[0][0] else tf.add_paragraph()
        p.text = cat
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COR_PRIMARIA
        p.space_after = Pt(6)
        
        for tech in techs:
            p = tf.add_paragraph()
            p.text = f"• {tech}"
            p.font.size = Pt(14)
            p.font.color.rgb = COR_TEXTO
            p.level = 1
            p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(6)
    
    # ==================== SLIDE 13: CONCLUSÕES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "CONCLUSÕES"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PRINCIPAIS CONCLUSÕES"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    conclusoes = [
        "Arquitetura em 4 camadas mostrou-se eficaz para aplicações mobile de saúde",
        "Padrão offline-first garantiu disponibilidade contínua (objetivo alcançado)",
        "Taxa de cache de ~85% demonstra eficiência da estratégia implementada",
        "Separação de responsabilidades facilitou desenvolvimento e manutenção",
        "Type safety do TypeScript preveniu erros em tempo de execução",
        "Autenticação multicamadas garante proteção adequada de dados sensíveis"
    ]
    
    for conclusao in conclusoes:
        p = tf.add_paragraph()
        p.text = f"✓ {conclusao}"
        p.font.size = Pt(15)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(7)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "CONTRIBUIÇÕES DO TRABALHO"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(10)
    
    contribuicoes = [
        "Arquitetura documentada e replicável para aplicações similares",
        "Implementação prática de offline-first em React Native",
        "Referência para desenvolvimento de apps acessíveis para idosos"
    ]
    
    for contrib in contribuicoes:
        p = tf.add_paragraph()
        p.text = f"• {contrib}"
        p.font.size = Pt(15)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    # ==================== SLIDE 14: LIMITAÇÕES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "LIMITAÇÕES"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "LIMITAÇÕES IDENTIFICADAS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(12)
    
    limitacoes = [
        ("Escopo de Testes", 
         "Testes realizados principalmente em ambiente de desenvolvimento, necessário validação com usuários reais idosos"),
        
        ("Plataforma", 
         "Versão atual focada em Android, build iOS requer macOS e Apple Developer Account"),
        
        ("Sincronização em Tempo Real", 
         "Implementação atual não possui sync bidirecional automática, requer refresh manual"),
        
        ("Catálogo de Exercícios", 
         "Base inicial de 10+ exercícios, expansão futura necessária para maior variedade"),
        
        ("Analytics e Monitoramento", 
         "Ausência de telemetria para rastreamento de uso e comportamento do usuário"),
        
        ("Internacionalização", 
         "Interface disponível apenas em português brasileiro")
    ]
    
    for titulo, desc in limitacoes:
        p = tf.add_paragraph()
        p.text = f"• {titulo}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COR_PRIMARIA
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COR_TEXTO
        p.level = 2
        p.space_after = Pt(8)
    
    # ==================== SLIDE 15: RECOMENDAÇÕES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = "RECOMENDAÇÕES E TRABALHOS FUTUROS"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "RECOMENDAÇÕES PARA TRABALHOS FUTUROS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_after = Pt(12)
    
    recomendacoes = [
        ("Sincronização em Tempo Real", 
         "Implementar WebSockets para sync bidirecional automática de dados"),
        
        ("Sistema de Fila Offline", 
         "Queue para operações pendentes quando offline, com retry automático"),
        
        ("Testes E2E Automatizados", 
         "Suite de testes end-to-end com Detox ou Maestro"),
        
        ("Push Notifications", 
         "Lembretes de treino e notificações de progresso para engajamento"),
        
        ("Analytics Detalhado", 
         "Implementar Firebase Analytics ou Amplitude para insights de uso"),
        
        ("Gamificação", 
         "Sistema de conquistas, badges e desafios para motivação"),
        
        ("Suporte Multiplataforma", 
         "Build para iOS e versão web progressive (PWA)"),
        
        ("Integração com Wearables", 
         "Conectar com smartwatches e monitores de atividade"),
        
        ("Inteligência Artificial", 
         "Recomendações personalizadas baseadas em ML")
    ]
    
    for titulo, desc in recomendacoes:
        p = tf.add_paragraph()
        p.text = f"→ {titulo}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COR_PRIMARIA
        p.space_after = Pt(2)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COR_TEXTO
        p.level = 1
        p.space_after = Pt(6)
    
    # ==================== SLIDE 16: AGRADECIMENTOS E ENCERRAMENTO ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "OBRIGADO!"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = COR_PRIMARIA
    title_p.alignment = PP_ALIGN.CENTER
    
    # Contato
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Guilherme Antony\n\nPerguntascontextual?"
    for p in contact_frame.paragraphs:
        p.font.size = Pt(24)
        p.font.color.rgb = COR_SECUNDARIA
        p.alignment = PP_ALIGN.CENTER
    
    # Download
    download_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    download_frame = download_box.text_frame
    download_frame.text = "📱 Download: https://expo.dev/artifacts/eas/s8rCmPjY3mcTXJXSYWtgqx.apk"
    download_p = download_frame.paragraphs[0]
    download_p.font.size = Pt(14)
    download_p.font.color.rgb = COR_SECUNDARIA
    download_p.alignment = PP_ALIGN.CENTER
    
    # Salvar apresentação
    output_path = '/home/antony/Documentos/Vivafit-Senior/Apresentacao_TCC_VivaFit_Seniors.pptx'
    prs.save(output_path)
    print(f"✓ Apresentação criada com sucesso: {output_path}")
    print(f"✓ Total de slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    criar_apresentacao()
